import os
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import hashlib
import json

# Document processing imports
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
import openpyxl

from .config import config

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.config = config
        self.documents_path = self.config.get_paths().get('documents', './data/documents')
        self.processed_path = self.config.get_paths().get('processed', './data/processed')
        self.chunk_size = self.config.get_vector_db_config().get('chunk_size', 1000)
        self.chunk_overlap = self.config.get_vector_db_config().get('chunk_overlap', 200)
        
        # Create directories if they don't exist
        Path(self.documents_path).mkdir(parents=True, exist_ok=True)
        Path(self.processed_path).mkdir(parents=True, exist_ok=True)
        
        # Load processed files tracking
        self.processed_files = self._load_processed_files()
    
    def _load_processed_files(self) -> Dict[str, str]:
        """Load the list of already processed files."""
        tracking_file = os.path.join(self.processed_path, 'processed_files.json')
        if os.path.exists(tracking_file):
            try:
                with open(tracking_file, 'r') as f:
                    return json.load(f)
            except Exception as e:
                logger.warning(f"Could not load processed files tracking: {e}")
        return {}
    
    def _save_processed_files(self):
        """Save the list of processed files."""
        tracking_file = os.path.join(self.processed_path, 'processed_files.json')
        try:
            with open(tracking_file, 'w') as f:
                json.dump(self.processed_files, f)
        except Exception as e:
            logger.error(f"Could not save processed files tracking: {e}")
    
    def _get_file_hash(self, file_path: str) -> str:
        """Generate hash for file content."""
        with open(file_path, 'rb') as f:
            return hashlib.md5(f.read()).hexdigest()
    
    def _is_file_processed(self, file_path: str) -> bool:
        """Check if file has already been processed."""
        file_hash = self._get_file_hash(file_path)
        return self.processed_files.get(file_path, '') == file_hash
    
    def _mark_file_processed(self, file_path: str):
        """Mark file as processed."""
        file_hash = self._get_file_hash(file_path)
        self.processed_files[file_path] = file_hash
        self._save_processed_files()
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            return ""
    
    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {e}")
            return ""
    
    def extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error extracting text from TXT {file_path}: {e}")
            return ""
    
    def extract_text_from_csv(self, file_path: str) -> str:
        """Extract text from CSV file."""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Error extracting text from CSV {file_path}: {e}")
            return ""
    
    def extract_text_from_xlsx(self, file_path: str) -> str:
        """Extract text from Excel file."""
        try:
            import pandas as pd
            df = pd.read_excel(file_path, sheet_name=None)
            text_parts = []
            for sheet_name, sheet_df in df.items():
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(sheet_df.to_string(index=False))
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from Excel file {file_path}: {e}")
            return ""

    def extract_text_from_doc(self, file_path: str) -> str:
        """Extract text from .doc file."""
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(file_path)
            text = doc.Content.Text
            doc.Close()
            word.Quit()
            return text
        except Exception as e:
            logger.error(f"Error extracting text from .doc file {file_path}: {e}")
            return ""

    def extract_text_from_ppt(self, file_path: str) -> str:
        """Extract text from .ppt file."""
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(file_path)
            text_parts = []
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if hasattr(shape, "TextFrame"):
                        text_parts.append(shape.TextFrame.TextRange.Text)
            presentation.Close()
            powerpoint.Quit()
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .ppt file {file_path}: {e}")
            return ""

    def extract_text_from_odt(self, file_path: str) -> str:
        """Extract text from OpenDocument Text file."""
        try:
            from zipfile import ZipFile
            with ZipFile(file_path) as zip_file:
                content = zip_file.read('content.xml')
                import xml.etree.ElementTree as ET
                root = ET.fromstring(content)
                text_elements = root.findall('.//{urn:oasis:names:tc:opendocument:xmlns:text:1.0}p')
                text_parts = [elem.text for elem in text_elements if elem.text]
                return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .odt file {file_path}: {e}")
            return ""

    def extract_text_from_ods(self, file_path: str) -> str:
        """Extract text from OpenDocument Spreadsheet file."""
        try:
            import pandas as pd
            df = pd.read_excel(file_path, sheet_name=None, engine='odf')
            text_parts = []
            for sheet_name, sheet_df in df.items():
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(sheet_df.to_string(index=False))
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .ods file {file_path}: {e}")
            return ""

    def extract_text_from_odp(self, file_path: str) -> str:
        """Extract text from OpenDocument Presentation file."""
        try:
            from zipfile import ZipFile
            with ZipFile(file_path) as zip_file:
                content = zip_file.read('content.xml')
                import xml.etree.ElementTree as ET
                root = ET.fromstring(content)
                text_elements = root.findall('.//{urn:oasis:names:tc:opendocument:xmlns:text:1.0}p')
                text_parts = [elem.text for elem in text_elements if elem.text]
                return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .odp file {file_path}: {e}")
            return ""

    def extract_text_from_rtf(self, file_path: str) -> str:
        """Extract text from RTF file."""
        try:
            import striprtf
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            return striprtf.rtf_to_text(rtf_content)
        except Exception as e:
            logger.error(f"Error extracting text from .rtf file {file_path}: {e}")
            return ""

    def extract_text_from_html(self, file_path: str) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            soup = BeautifulSoup(html_content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            logger.error(f"Error extracting text from .html file {file_path}: {e}")
            return ""

    def extract_text_from_xls(self, file_path: str) -> str:
        """Extract text from .xls file."""
        try:
            import pandas as pd
            df = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
            text_parts = []
            for sheet_name, sheet_df in df.items():
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(sheet_df.to_string(index=False))
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .xls file {file_path}: {e}")
            return ""

    def extract_text_from_eml(self, file_path: str) -> str:
        """Extract text from .eml file."""
        try:
            import email
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                msg = email.message_from_file(f)
            text_parts = []
            if msg['subject']:
                text_parts.append(f"Subject: {msg['subject']}")
            if msg['from']:
                text_parts.append(f"From: {msg['from']}")
            if msg['to']:
                text_parts.append(f"To: {msg['to']}")
            text_parts.append("Body:")
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    text_parts.append(part.get_payload(decode=True).decode())
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .eml file {file_path}: {e}")
            return ""

    def extract_text_from_msg(self, file_path: str) -> str:
        """Extract text from .msg file."""
        try:
            import win32com.client
            outlook = win32com.client.Dispatch("Outlook.Application")
            msg = outlook.CreateItemFromTemplate(file_path)
            text_parts = []
            if msg.Subject:
                text_parts.append(f"Subject: {msg.Subject}")
            if msg.SenderName:
                text_parts.append(f"From: {msg.SenderName}")
            if msg.Body:
                text_parts.append(f"Body: {msg.Body}")
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .msg file {file_path}: {e}")
            return ""

    def extract_text_from_archive(self, file_path: str) -> str:
        """Extract text from archive files."""
        try:
            import tempfile
            import zipfile
            import tarfile
            import os
            
            extracted_texts = []
            temp_dir = tempfile.mkdtemp()
            
            try:
                if file_path.endswith('.zip'):
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        zip_ref.extractall(temp_dir)
                elif file_path.endswith('.tar') or file_path.endswith('.tar.gz') or file_path.endswith('.gz'):
                    with tarfile.open(file_path, 'r:*') as tar_ref:
                        tar_ref.extractall(temp_dir)
                elif file_path.endswith('.rar'):
                    import rarfile
                    with rarfile.RarFile(file_path, 'r') as rar_ref:
                        rar_ref.extractall(temp_dir)
                elif file_path.endswith('.7z'):
                    import py7zr
                    with py7zr.SevenZipFile(file_path, 'r') as sz_ref:
                        sz_ref.extractall(temp_dir)
                
                # Process extracted files
                for root, dirs, files in os.walk(temp_dir):
                    for file in files:
                        file_path_extracted = os.path.join(root, file)
                        try:
                            text = self.extract_text_from_file(file_path_extracted)
                            if text:
                                extracted_texts.append(f"File: {file}\n{text}")
                        except:
                            continue
                            
            finally:
                import shutil
                shutil.rmtree(temp_dir)
            
            return "\n\n".join(extracted_texts)
        except Exception as e:
            logger.error(f"Error extracting text from archive {file_path}: {e}")
            return ""

    def extract_text_from_epub(self, file_path: str) -> str:
        """Extract text from EPUB file."""
        try:
            from zipfile import ZipFile
            import xml.etree.ElementTree as ET
            
            text_parts = []
            with ZipFile(file_path) as zip_file:
                # Read container.xml to find content files
                container = zip_file.read('META-INF/container.xml')
                root = ET.fromstring(container)
                rootfile = root.find('.//{urn:oasis:names:tc:opendocument:xmlns:container}rootfile')
                if rootfile is not None:
                    content_path = rootfile.get('full-path')
                    content = zip_file.read(content_path)
                    content_root = ET.fromstring(content)
                    
                    # Extract text from content
                    for elem in content_root.iter():
                        if elem.text and elem.text.strip():
                            text_parts.append(elem.text.strip())
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from .epub file {file_path}: {e}")
            return ""

    def extract_text_from_mobi(self, file_path: str) -> str:
        """Extract text from MOBI file."""
        try:
            # MOBI files are complex binary formats, this is a basic implementation
            with open(file_path, 'rb') as f:
                content = f.read()
            # Look for text content in the MOBI file
            # This is a simplified approach - full MOBI parsing would require a dedicated library
            text_content = content.decode('utf-8', errors='ignore')
            # Remove binary data and keep only printable text
            import re
            text_content = re.sub(r'[^\x20-\x7E\n\r\t]', '', text_content)
            return text_content
        except Exception as e:
            logger.error(f"Error extracting text from .mobi file {file_path}: {e}")
            return ""

    def extract_text_from_azw3(self, file_path: str) -> str:
        """Extract text from AZW3 file."""
        try:
            # AZW3 is Amazon's format, similar to MOBI
            with open(file_path, 'rb') as f:
                content = f.read()
            text_content = content.decode('utf-8', errors='ignore')
            import re
            text_content = re.sub(r'[^\x20-\x7E\n\r\t]', '', text_content)
            return text_content
        except Exception as e:
            logger.error(f"Error extracting text from .azw3 file {file_path}: {e}")
            return ""
    
    def extract_text_from_file(self, file_path: str) -> str:
        """Extract text from file based on its extension."""
        file_ext = Path(file_path).suffix.lower()
        
        extractors = {
            # Document formats
            '.pdf': self.extract_text_from_pdf,
            '.docx': self.extract_text_from_docx,
            '.doc': self.extract_text_from_doc,
            '.pptx': self.extract_text_from_pptx,
            '.ppt': self.extract_text_from_ppt,
            '.odt': self.extract_text_from_odt,
            '.ods': self.extract_text_from_ods,
            '.odp': self.extract_text_from_odp,
            
            # Text formats
            '.txt': self.extract_text_from_txt,
            '.md': self.extract_text_from_txt,
            '.rtf': self.extract_text_from_rtf,
            '.tex': self.extract_text_from_txt,
            '.log': self.extract_text_from_txt,
            '.ini': self.extract_text_from_txt,
            '.cfg': self.extract_text_from_txt,
            '.conf': self.extract_text_from_txt,
            '.json': self.extract_text_from_txt,
            '.xml': self.extract_text_from_txt,
            '.html': self.extract_text_from_html,
            '.htm': self.extract_text_from_html,
            '.css': self.extract_text_from_txt,
            '.js': self.extract_text_from_txt,
            '.py': self.extract_text_from_txt,
            '.java': self.extract_text_from_txt,
            '.cpp': self.extract_text_from_txt,
            '.c': self.extract_text_from_txt,
            '.h': self.extract_text_from_txt,
            '.cs': self.extract_text_from_txt,
            '.php': self.extract_text_from_txt,
            '.rb': self.extract_text_from_txt,
            '.go': self.extract_text_from_txt,
            '.rs': self.extract_text_from_txt,
            '.swift': self.extract_text_from_txt,
            '.kt': self.extract_text_from_txt,
            '.scala': self.extract_text_from_txt,
            '.sql': self.extract_text_from_txt,
            '.sh': self.extract_text_from_txt,
            '.bat': self.extract_text_from_txt,
            '.ps1': self.extract_text_from_txt,
            '.yaml': self.extract_text_from_txt,
            '.yml': self.extract_text_from_txt,
            '.toml': self.extract_text_from_txt,
            '.csv': self.extract_text_from_csv,
            '.tsv': self.extract_text_from_csv,
            
            # Spreadsheet formats
            '.xlsx': self.extract_text_from_xlsx,
            '.xls': self.extract_text_from_xls,
            '.xlsm': self.extract_text_from_xlsx,
            '.xlsb': self.extract_text_from_xlsx,
            
            # Email formats
            '.eml': self.extract_text_from_eml,
            '.msg': self.extract_text_from_msg,
            
            # Archive formats (extract and process)
            '.zip': self.extract_text_from_archive,
            '.rar': self.extract_text_from_archive,
            '.7z': self.extract_text_from_archive,
            '.tar': self.extract_text_from_archive,
            '.gz': self.extract_text_from_archive,
            
            # Other formats
            '.epub': self.extract_text_from_epub,
            '.mobi': self.extract_text_from_mobi,
            '.azw3': self.extract_text_from_azw3,
        }
        
        if file_ext in extractors:
            try:
                return extractors[file_ext](file_path)
            except Exception as e:
                logger.error(f"Error extracting text from {file_path}: {e}")
                return ""
        else:
            # Try to extract as plain text for unknown file types
            logger.info(f"Attempting to extract text from unknown file type: {file_ext}")
            return self.extract_text_from_txt(file_path)
    
    def chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks."""
        if not text.strip():
            return []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + self.chunk_size
            
            # If this is not the first chunk, include overlap
            if start > 0:
                start = start - self.chunk_overlap
            
            # Extract the chunk
            chunk = text[start:end].strip()
            
            if chunk:
                chunks.append(chunk)
            
            # Move to next chunk
            start = end
            
            # If we've reached the end, break
            if end >= len(text):
                break
        
        return chunks
    
    def process_documents(self, force: bool = False) -> List[Dict[str, Any]]:
        """
        Process all documents in the documents folder.
        
        Args:
            force (bool): If True, re-process all documents even if already processed
            
        Returns:
            List[Dict[str, Any]]: List of processed document chunks
        """
        processed_chunks = []
        
        if not os.path.exists(self.documents_path):
            logger.warning(f"Documents path does not exist: {self.documents_path}")
            return processed_chunks
        
        # Define all supported file extensions
        supported_extensions = {
            # Document formats
            '.pdf', '.docx', '.doc', '.pptx', '.ppt', '.odt', '.ods', '.odp',
            # Text formats
            '.txt', '.md', '.rtf', '.tex', '.log', '.ini', '.cfg', '.conf', 
            '.json', '.xml', '.html', '.htm', '.css', '.js', '.py', '.java', 
            '.cpp', '.c', '.h', '.cs', '.php', '.rb', '.go', '.rs', '.swift', 
            '.kt', '.scala', '.sql', '.sh', '.bat', '.ps1', '.yaml', '.yml', 
            '.toml', '.csv', '.tsv',
            # Spreadsheet formats
            '.xlsx', '.xls', '.xlsm', '.xlsb',
            # Email formats
            '.eml', '.msg',
            # Archive formats
            '.zip', '.rar', '.7z', '.tar', '.gz',
            # Other formats
            '.epub', '.mobi', '.azw3'
        }
        
        for file_path in Path(self.documents_path).rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                file_path_str = str(file_path)
                
                # Check if file has already been processed (unless force=True)
                if not force and self._is_file_processed(file_path_str):
                    logger.info(f"File already processed: {file_path.name}")
                    continue
                
                logger.info(f"Processing file: {file_path.name}")
                
                # Extract text
                text = self.extract_text_from_file(file_path_str)
                
                if text:
                    # Chunk the text
                    chunks = self.chunk_text(text)
                    
                    # Get directory structure information
                    dir_structure = self._get_directory_structure(file_path)
                    
                    # Create chunk metadata with enhanced information
                    for i, chunk in enumerate(chunks):
                        chunk_data = {
                            'text': chunk,
                            'file_path': file_path_str,
                            'file_name': file_path.name,
                            'chunk_index': i,
                            'total_chunks': len(chunks),
                            'file_type': file_path.suffix.lower(),
                            'directory_structure': dir_structure,
                            'relative_path': str(file_path.relative_to(Path(self.documents_path))),
                            'parent_directories': self._get_parent_directories(file_path),
                            'file_size': os.path.getsize(file_path_str),
                            'last_modified': os.path.getmtime(file_path_str)
                        }
                        processed_chunks.append(chunk_data)
                    
                    # Mark file as processed
                    self._mark_file_processed(file_path_str)
                    logger.info(f"Processed {len(chunks)} chunks from {file_path.name}")
                else:
                    logger.warning(f"No text extracted from {file_path.name}")
        
        return processed_chunks
    
    def _get_directory_structure(self, file_path: Path) -> Dict[str, Any]:
        """Get the directory structure information for a file."""
        try:
            # Get the relative path from documents directory
            relative_path = file_path.relative_to(Path(self.documents_path))
            path_parts = relative_path.parts
            
            # Analyze directory structure
            structure = {
                'full_path': str(relative_path),
                'path_parts': path_parts,
                'depth': len(path_parts) - 1,  # -1 because last part is filename
                'parent_directories': path_parts[:-1] if len(path_parts) > 1 else [],
                'file_name': path_parts[-1] if path_parts else file_path.name,
                'directory_hierarchy': self._build_hierarchy(path_parts)
            }
            
            # Add semantic analysis for common patterns
            structure['semantic_info'] = self._analyze_semantic_structure(path_parts)
            
            return structure
            
        except Exception as e:
            logger.error(f"Error getting directory structure for {file_path}: {e}")
            return {
                'full_path': str(file_path),
                'path_parts': [file_path.name],
                'depth': 0,
                'parent_directories': [],
                'file_name': file_path.name,
                'directory_hierarchy': {},
                'semantic_info': {}
            }
    
    def _get_parent_directories(self, file_path: Path) -> List[str]:
        """Get list of parent directories for a file."""
        try:
            relative_path = file_path.relative_to(Path(self.documents_path))
            path_parts = relative_path.parts
            return list(path_parts[:-1]) if len(path_parts) > 1 else []
        except Exception:
            return []
    
    def _build_hierarchy(self, path_parts: tuple) -> Dict[str, Any]:
        """Build a hierarchical structure from path parts."""
        hierarchy = {}
        current = hierarchy
        
        for part in path_parts[:-1]:  # Exclude filename
            current[part] = {'type': 'directory', 'children': {}}
            current = current[part]['children']
        
        if path_parts:
            current[path_parts[-1]] = {'type': 'file'}
        
        return hierarchy
    
    def _analyze_semantic_structure(self, path_parts: tuple) -> Dict[str, Any]:
        """Analyze semantic meaning of directory structure."""
        semantic_info = {}
        
        # Common patterns for academic/course materials
        course_patterns = {
            'module': ['module', 'mod', 'unit', 'chapter', 'week'],
            'subject': ['python', 'java', 'cpp', 'c++', 'database', 'dbms', 'ml', 'ai', 'ds', 'algorithms'],
            'type': ['notes', 'syllabus', 'assignment', 'lab', 'tutorial', 'lecture', 'slides', 'exam'],
            'level': ['basic', 'intermediate', 'advanced', 'beginner', 'expert']
        }
        
        path_text = ' '.join(path_parts).lower()
        
        # Detect patterns
        for category, patterns in course_patterns.items():
            matches = []
            for pattern in patterns:
                if pattern in path_text:
                    matches.append(pattern)
            if matches:
                semantic_info[category] = matches
        
        # Detect module numbers
        import re
        module_numbers = re.findall(r'module\s*(\d+)|mod\s*(\d+)|unit\s*(\d+)|week\s*(\d+)', path_text, re.IGNORECASE)
        if module_numbers:
            semantic_info['module_numbers'] = [num for group in module_numbers for num in group if num]
        
        # Detect file types and purposes
        file_types = {
            'notes': ['notes', 'note', 'lecture', 'tutorial'],
            'syllabus': ['syllabus', 'curriculum', 'outline'],
            'assignment': ['assignment', 'homework', 'project', 'lab'],
            'slides': ['slides', 'presentation', 'ppt', 'pptx'],
            'exam': ['exam', 'test', 'quiz', 'question']
        }
        
        for file_type, keywords in file_types.items():
            if any(keyword in path_text for keyword in keywords):
                semantic_info['file_purpose'] = file_type
                break
        
        return semantic_info
    
    def process_single_file(self, file_path: str) -> List[Dict[str, Any]]:
        """Process a single file."""
        if not os.path.exists(file_path):
            logger.error(f"File does not exist: {file_path}")
            return []
        
        # Check if file has already been processed
        if self._is_file_processed(file_path):
            logger.info(f"File already processed: {file_path}")
            return []
        
        logger.info(f"Processing file: {file_path}")
        
        # Extract text
        text = self.extract_text_from_file(file_path)
        
        if not text:
            logger.warning(f"No text extracted from {file_path}")
            return []
        
        # Chunk the text
        chunks = self.chunk_text(text)
        
        # Create chunk metadata
        processed_chunks = []
        for i, chunk in enumerate(chunks):
            chunk_data = {
                'text': chunk,
                'file_path': file_path,
                'file_name': Path(file_path).name,
                'chunk_index': i,
                'total_chunks': len(chunks),
                'file_type': Path(file_path).suffix.lower()
            }
            processed_chunks.append(chunk_data)
        
        # Mark file as processed
        self._mark_file_processed(file_path)
        logger.info(f"Processed {len(chunks)} chunks from {file_path}")
        
        return processed_chunks
